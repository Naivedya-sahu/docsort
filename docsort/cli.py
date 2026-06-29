#!/usr/bin/env python3
"""
docsort.cli — tag academic documents with a local (or frontier) LLM, then move by prefix.

Tag vocabulary lives in TAGS.md (single source) and is injected into the model prompt.
Classification tiers: TEXT -> 5-page ESCALATE -> VISION (+page-3) -> FILENAME.
On a hard 99UNS, an optional FRONTIER model (Claude Code CLI, haiku) gives a verdict.

  docsort "D:\\AcademicsCOPY"                       # dry-run (config supplies host/model)
  docsort "D:\\AcademicsCOPY" --apply               # rename in place
  docsort "D:\\AcademicsCOPY" --copy --apply        # tag a copy, originals untouched
  docsort "D:\\AcademicsCOPY" --frontier claude      # Claude Code (haiku) for hard cases
  docsort "D:\\AcademicsCOPY" --move "D:\\Archive" --apply
  docsort --edit-tags                                # open your TAGS.md in an editor

Backends: --backend local (main) ; --frontier none|claude|cmd (99UNS fallback, claude=haiku).
Deps: see requirements.txt.
"""
from __future__ import annotations
import os, sys, json, csv, re, time, base64, argparse, subprocess, shutil, shlex, urllib.request

try:                                   # works installed (package) and as `python -m docsort.cli`
    from .config import (load_config, resolve_api, resolve_location, arg_defaults,
                         tags_path, prompt_path, config_path, user_dir)
except ImportError:                    # fallback if run as a loose script
    from config import (load_config, resolve_api, resolve_location, arg_defaults,
                        tags_path, prompt_path, config_path, user_dir)

def unique_path(p):
    """Avoid overwriting: foo.pdf -> foo__1.pdf if taken."""
    if not os.path.exists(p): return p
    base,ext=os.path.splitext(p); i=1
    while os.path.exists(f"{base}__{i}{ext}"): i+=1
    return f"{base}__{i}{ext}"

def make_working_copy(root):
    """Copy root -> <root>COPY (unique) and return the copy path, so originals stay untouched."""
    base=root.rstrip("/\\"); dst=base+"COPY"; i=1
    while os.path.exists(dst): dst=f"{base}COPY{i}"; i+=1
    shutil.copytree(root,dst); return dst

MISC_DIR="misc"
def move_to_misc(root,cur):
    """Quarantine an unsure (99UNS) file into <root>/misc/. Returns the new path."""
    d=os.path.join(root,MISC_DIR); os.makedirs(d,exist_ok=True)
    tgt=unique_path(os.path.join(d,os.path.basename(cur))); shutil.move(cur,tgt); return tgt

def edit_file(p):
    """Open a file in the OS default editor (Notepad on Windows)."""
    print("opening for edit:",p)
    try:
        if os.name=="nt": os.startfile(p)                       # type: ignore[attr-defined]
        else: subprocess.Popen([os.environ.get("EDITOR","nano"),p])
    except Exception as e:
        print("could not open an editor:",e,"\nedit this file manually:",p)

EXT_TEXT={".pdf",".txt",".md",".docx",".pptx",".ppt",".doc"}
MIN_TEXT=80; DEEP_PAGES=5; DEEP_CAP=4000; DPI=120
STREAMS=set(); SUBJECTS=set(); TYPES=set()   # filled from TAGS.md
_STATS={"ttoks":0,"ok":0,"empty":0}          # per-file model accounting (reset each file)

def journal_done(root):
    """rel-paths already 'done' in the run journal -> (rel, mtime) set, for --resume."""
    p=os.path.join(root,"_docsort_state.jsonl"); seen={}
    if not os.path.isfile(p): return seen
    for ln in open(p,encoding="utf-8",errors="replace"):
        try:
            j=json.loads(ln)
            if j.get("status")=="done": seen[j.get("rel","")]=j.get("mtime",0)
        except Exception: pass
    return seen

# ---------- tag loading (single source) ----------
def load_tags(path):
    streams={}; subjects={}; types=[]
    try: txt=open(path,encoding="utf-8").read()
    except Exception: txt=""
    def block(header):
        m=re.search(r"##\s*"+header+r".*?```tags\n(.*?)```",txt,re.S|re.I)
        return m.group(1).strip().splitlines() if m else []
    for ln in block("STREAMS"):
        if ln.strip(): streams[ln.split()[0]]=ln.strip()
    for ln in block("SUBJECTS"):
        if ln.strip(): subjects[ln.split()[0]]=ln.strip()
    m=re.search(r"##\s*TYPES.*?```tags\n(.*?)```",txt,re.S|re.I)
    if m: types=[x.split()[0] for x in m.group(1).splitlines() if x.strip()]
    if not streams: streams={"CW":"CW","GATE":"GATE","PROJ":"PROJ","RES":"RES","REC":"REC","REF":"REF"}
    if not subjects: subjects={"99UNS":"unsure","NA":"na"}
    if not types: types=["misc"]
    return streams,subjects,types

def build_system(template_path, streams, subjects, types):
    try:
        s=open(template_path,encoding="utf-8").read()
        m=re.search(r"<SYSTEM>(.*?)</SYSTEM>",s,re.S); tmpl=m.group(1) if m else s
    except Exception:
        tmpl="Reply: STREAM SUBJECT TYPE CONF.\n{{STREAMS}}\n{{SUBJECTS}}\n{{TYPES}}"
    tmpl=tmpl.replace("{{STREAMS}}","\n".join("  "+v for v in streams.values()))
    tmpl=tmpl.replace("{{SUBJECTS}}","\n".join("  "+v for v in subjects.values()))
    tmpl=tmpl.replace("{{TYPES}}"," ".join(types))
    return tmpl.strip()

# ---------- text / image extraction ----------
def pdf_text(path,pages=2,cap=2000):
    try:
        import fitz; d=fitz.open(path); t=""
        for i in range(min(pages,d.page_count)):
            t+=d[i].get_text()
            if len(t)>cap: break
        d.close(); return t[:cap]
    except Exception: return ""
_WORD=None
def _doc_legacy(path):
    """Read old binary .doc via a reused Word COM instance (Windows + Word + pywin32)."""
    global _WORD
    try:
        if _WORD is None:
            import win32com.client as wc
            _WORD=wc.Dispatch("Word.Application"); _WORD.Visible=False
        d=_WORD.Documents.Open(os.path.abspath(path), ReadOnly=True, AddToRecentFiles=False)
        t=d.Content.Text; d.Close(False); return t
    except Exception: return ""   # no Word/pywin32 -> falls back to filename tier
def _word_quit():
    global _WORD
    try:
        if _WORD is not None: _WORD.Quit(); _WORD=None
    except Exception: pass

def doc_text(path,cap=2000):
    e=os.path.splitext(path)[1].lower()
    try:
        if e==".pdf": return pdf_text(path,2,cap)
        if e in (".txt",".md"): return open(path,encoding="utf-8",errors="replace").read()[:cap]
        if e==".doc": return _doc_legacy(path)[:cap]
        if e==".docx":
            import docx; return "\n".join(p.text for p in docx.Document(path).paragraphs)[:cap]
        if e==".pptx":
            from pptx import Presentation
            return " ".join(s.text for sl in Presentation(path).slides for s in sl.shapes
                             if getattr(s,'has_text_frame',False))[:cap]
    except Exception: return ""
    return ""
def page_png(path,page=0,dpi=None):
    try:
        import fitz; d=fitz.open(path)
        if d.page_count==0: return None
        if d.page_count<=page: page=0
        b=d[page].get_pixmap(dpi=dpi or DPI).tobytes("png"); d.close(); return b
    except Exception: return None

# ---------- backends ----------
def _http(url,payload,headers):
    req=urllib.request.Request(url,data=json.dumps(payload).encode(),headers=headers)
    return json.load(urllib.request.urlopen(req,timeout=180))

def available_models(api):
    try:
        r=json.load(urllib.request.urlopen(api.replace("/chat/completions","/models"),timeout=8))
        return [m.get("id") for m in r.get("data",[])]
    except Exception: return []

def resolve_model(api,want,prefer_vision=False):
    """If 'want' isn't loaded on the server, fall back to a loaded model (vision-preferring).
    Returns (model_id, server_reachable). Works for any user: picks whatever VL model is loaded."""
    av=available_models(api)
    if not av: return want, False
    if want in av: return want, True
    pick=None
    if prefer_vision:
        pick=next((m for m in av if ("vl" in m.lower() or "vision" in m.lower())),None)
    pick=pick or next((m for m in av if "embed" not in m.lower()), av[0])
    return pick, True

def llm(a,backend,system,user_text,image_png=None):
    """Return raw model text. backend: local | claude | cmd."""
    if backend=="local":
        if image_png:
            content=[{"type":"text","text":user_text},
                     {"type":"image_url","image_url":{"url":"data:image/png;base64,"+base64.b64encode(image_png).decode()}}]
        else:
            content=user_text
        msgs=[{"role":"system","content":system},{"role":"user","content":content}]
        try:
            model=a.vision_model if image_png else a.model
            r=_http(a.api,{"model":model,"temperature":0,"max_tokens":24,"messages":msgs},
                    {"Content-Type":"application/json"})
            u=r.get("usage") or {}; _STATS["ttoks"]+=int(u.get("total_tokens",0) or u.get("completion_tokens",0)); _STATS["ok"]+=1
            return r["choices"][0]["message"]["content"]
        except Exception:
            _STATS["empty"]+=1
            return ""   # -> decide() -> 99UNS; one bad call never kills the run
    if backend in ("claude","cmd"):   # subscription CLIs — NO API key. Text only.
        import tempfile
        oneshot=(system+"\n\n"+user_text+
                 "\n\nRespond with ONLY one line: STREAM SUBJECT TYPE CONF. "
                 "No explanation, no markdown, no questions.")
        neutral=tempfile.gettempdir()   # run outside the repo so the CLI doesn't read project files
        try:
            if backend=="claude":       # Claude Code CLI -> Claude subscription
                exe=shutil.which("claude") or "claude"
                out=subprocess.run([exe,"-p","--model","haiku",oneshot],capture_output=True,text=True,
                                   timeout=180,cwd=neutral)   # bound to haiku (sub-covered, standard context); resolved .exe -> no shell
                return out.stdout
            tmpl=getattr(a,"frontier_cmd","") or ""   # cmd: templated CLI from config, prompt via stdin
            if not tmpl: return ""
            out=subprocess.run(shlex.split(tmpl),input=oneshot,capture_output=True,text=True,
                               timeout=180,cwd=neutral)   # no shell -> no injection from the template
            return out.stdout
        except Exception: return ""
    return ""

def parse(out):
    o=(out or "").upper()
    st=next((s for s in sorted(STREAMS,key=len,reverse=True) if re.search(r'\b'+re.escape(s)+r'\b',o)),"CW")
    su=next((c for c in sorted(SUBJECTS,key=len,reverse=True) if re.search(r'\b'+re.escape(c)+r'\b',o)),"99UNS")
    ty=next((t for t in TYPES if re.search(r'\b'+re.escape(t.upper())+r'\b',o)),"misc")
    cf="high" if "HIGH" in o else "low"
    m=re.search(r'PROPOSE[:\s]+([A-Z][A-Z0-9_]{1,11})',o)
    prop=m.group(1) if m else ""
    return st,su,ty,cf,prop

def decide(out):
    """parse + fold a model proposal into a ~LABEL subject (review symbol)."""
    st,su,ty,cf,prop=parse(out)
    if prop and su=="99UNS": su="~"+prop
    return st,su,ty,cf

def classify(a,sysp,full,fn,rel):
    ispdf=full.lower().endswith(".pdf")
    snip=doc_text(full)
    u=lambda txt:f"Filename: {fn}\nFolder: {rel}\nText:\n{txt[:DEEP_CAP]}\n\nAnswer (STREAM SUBJECT TYPE CONF):"
    if len(snip.strip())>=MIN_TEXT:
        st,su,ty,cf=decide(llm(a,a.backend,sysp,u(snip))); src="text"
        if su=="99UNS" and ispdf:
            deep=pdf_text(full,DEEP_PAGES,DEEP_CAP)
            if len(deep)>len(snip):
                r=decide(llm(a,a.backend,sysp,u(deep)))
                if r[1]!="99UNS": st,su,ty,cf,src=(*r,"text5"); snip=deep
        if su=="99UNS" and a.frontier!="none":
            r=decide(llm(a,a.frontier,sysp,u(snip)))
            if r[1]!="99UNS": return (*r,"frontier:"+a.frontier)
        return st,su,ty,cf,src
    if a.vision and ispdf and (png:=page_png(full,0)):
        un="(handwritten/scanned page) Answer (STREAM SUBJECT TYPE CONF):"
        st,su,ty,cf=decide(llm(a,a.backend,sysp,f"Filename: {fn}\nFolder: {rel}\n"+un,png)); src="vision"
        if su=="99UNS" and (p3:=page_png(full,2)):
            r=decide(llm(a,a.backend,sysp,f"Filename: {fn}\nFolder: {rel}\n(page 3) "+un,p3))
            if r[1]!="99UNS": return (*r,"vision3")
        return st,su,ty,cf,src
    r=decide(llm(a,a.backend,sysp,u("")))
    if r[1]=="99UNS" and a.frontier!="none":
        rf=decide(llm(a,a.frontier,sysp,u("")))
        if rf[1]!="99UNS": return (*rf,"frontier:"+a.frontier)
    return (*r,"filename")

def move_by_prefix(root,dest,apply):
    pat=re.compile(r'^\[([A-Z]+)-([0-9A-Z]+)\]\s+(.*)$'); n=0; rows=[]
    for dp,_,fns in os.walk(root):
        for fn in fns:
            m=pat.match(fn)
            if not m: continue
            src=os.path.join(dp,fn); tgt=os.path.join(dest,m.group(1),m.group(2))
            newp=os.path.join(tgt,m.group(3))
            if apply:
                os.makedirs(tgt,exist_ok=True)
                newp=unique_path(newp)                       # never overwrite
                try: shutil.move(src,newp)                   # cross-device safe
                except Exception as e: print("move-fail:",fn,e)
            rows.append((src,newp)); n+=1
    log=os.path.join(dest if apply else root,"_move_log.csv")
    try:
        os.makedirs(os.path.dirname(log),exist_ok=True)
        with open(log,"w",encoding="utf-8",newline="") as g: csv.writer(g).writerows([("from","to")]+rows)
    except Exception: pass
    print(f"{'MOVED' if apply else 'DRY-RUN move'}: {n} files -> {dest}")
    return n

SKIP_NAMES={"tag-review.md","readme.md","changelog.md","guide.md","troubleshooting.md","tags.md"}
def iter_targets(root):
    for dp,_,fns in os.walk(root):
        for fn in fns:
            if fn.lower() in SKIP_NAMES or fn.startswith(("_docsort","_doc_handler","_move")): continue
            if os.path.splitext(fn)[1].lower() in EXT_TEXT and not fn.startswith("["):
                yield dp,fn,fn                       # (dir, current name, base name)

def _folder_match(full, root, folder):
    """True if `full`'s directory is under `folder` (absolute prefix, or relative segment/prefix of root)."""
    d=os.path.dirname(os.path.abspath(full)).replace("\\","/")
    f=folder.replace("\\","/").rstrip("/")
    if not f: return False
    if os.path.isabs(folder): return d==f or d.startswith(f+"/")
    rel=os.path.relpath(os.path.dirname(full),root).replace("\\","/")
    rel="" if rel=="." else rel
    parts=rel.split("/") if rel else []
    return rel==f or rel.startswith(f+"/") or f in parts

def passes_filter(full, root, include, exclude):
    if exclude and any(_folder_match(full,root,x) for x in exclude): return False
    if include and not any(_folder_match(full,root,x) for x in include): return False
    return True

_TAGPFX=re.compile(r'^\[[^\]]+\]\s+(.*)$')
def iter_tagged(root):
    """Already-prefixed docs, for --retag. Yields (dir, current name, stripped base)."""
    for dp,_,fns in os.walk(root):
        for fn in fns:
            m=_TAGPFX.match(fn)
            if m and os.path.splitext(m.group(1))[1].lower() in EXT_TEXT:
                yield dp,fn,m.group(1)

def review(root, log=None):
    """Aggregate a run log into TAG-REVIEW.md: distribution, proposed (~) tags, low-conf. No model needed."""
    import collections
    logp=log or os.path.join(root,"_docsort_log.csv")
    if not os.path.isfile(logp): print("no log at",logp); return
    rows=list(csv.DictReader(open(logp,encoding="utf-8")))
    combo=collections.Counter(f"{r['stream']}-{r['subject']}" for r in rows)
    props=collections.Counter(r['subject'] for r in rows if r['subject'].startswith('~'))
    low=[r for r in rows if r['conf']=='low']
    out=["# TAG REVIEW","",f"{len(rows)} files · {len(props)} proposed tag(s) · {len(low)} low-confidence.","",
         "## Distribution (stream-subject)"]
    for k,n in combo.most_common(): out.append(f"- `{k}` x {n}")
    out.append("\n## (review) Proposed tags -> promote to TAGS.md")
    if props:
        for lab,n in props.most_common():
            out.append(f"- `{lab}` x {n}  -> if recurring, add a SUBJECT code for `{lab[1:]}` in TAGS.md, then re-run")
    else: out.append("- none")
    out.append("\n## Low-confidence (eyeball)")
    for r in low[:50]: out.append(f"- {r['stream']}-{r['subject']} · {r['source']} · {r['old']}")
    p=os.path.join(root,"TAG-REVIEW.md"); open(p,"w",encoding="utf-8").write("\n".join(out))
    print(f"review -> {p}")
    print(f"proposals: {dict(props) or 'none'} | low-conf: {len(low)} | groups: {len(combo)}")

def _journal_rows(root):
    """Dedup journal -> latest row per file."""
    p=os.path.join(root,"_docsort_state.jsonl")
    if not os.path.isfile(p): return []
    last={}
    for ln in open(p,encoding="utf-8",errors="replace"):
        try: j=json.loads(ln); last[j.get("rel")]=j
        except Exception: pass
    return list(last.values())

def report(root):
    """Build DOCSORT-REPORT.md from the journal + append a run summary to the global index."""
    import collections
    rows=_journal_rows(root)
    if not rows: print("no journal at",root); return None
    done=[r for r in rows if r.get("status")=="done"]
    fails=[r for r in rows if r.get("status")=="failed"]
    combo=collections.Counter(f"{r['stream']}-{r['subject']}" for r in done)
    types=collections.Counter(r.get("type") for r in done)
    props=collections.Counter(r["subject"] for r in done if str(r.get("subject","")).startswith("~"))
    low=[r for r in done if r.get("conf")=="low"]
    out=["# docsort report","",f"{len(rows)} files · {len(done)} done · {len(fails)} failed · "
         f"{len(props)} proposals · {len(low)} low-conf","","## stream-subject"]
    out+=[f"- `{k}` x{c}" for k,c in combo.most_common()]
    out+=["","## types"]+[f"- {k} x{c}" for k,c in types.most_common()]
    if props: out+=["","## proposals (~ — promote in Edit Tags, then --retag)"]+[f"- {k} x{c}" for k,c in props.most_common()]
    if low: out+=["","## low-confidence (eyeball)"]+[f"- {r['stream']}-{r['subject']} · {r.get('source')} · {r.get('name')}" for r in low[:50]]
    if fails: out+=["","## failed"]+[f"- {r.get('name')} — {r.get('error')}" for r in fails[:50]]
    rp=os.path.join(root,"DOCSORT-REPORT.md"); open(rp,"w",encoding="utf-8").write("\n".join(out))
    try:
        rec={"ts":int(time.time()),"root":root,"n":len(rows),"done":len(done),
             "failed":len(fails),"by":dict(combo)}
        open(os.path.join(user_dir(),"index.jsonl"),"a",encoding="utf-8").write(json.dumps(rec)+"\n")
    except Exception: pass
    print(f"report -> {rp}  ({len(done)} done, {len(fails)} failed)")
    return rp

def undo(root):
    """Reverse renames/moves recorded in the journal (restore originals)."""
    moves=[(j["dst"],j["rel"]) for j in _journal_rows(root)
           if j.get("status")=="done" and j.get("dst") and j.get("dst")!=j.get("rel")]
    n=0
    for dst,rel in reversed(moves):
        src=os.path.join(root,dst); tgt=os.path.join(root,rel)
        if os.path.exists(src):
            try:
                os.makedirs(os.path.dirname(tgt),exist_ok=True); shutil.move(src,unique_path(tgt)); n+=1
            except Exception as e: print("  undo-fail:",dst,e)
    print(f"undo: restored {n} files in {root}")

def stats():
    """Lifetime totals from the global index."""
    import collections
    idx=os.path.join(user_dir(),"index.jsonl")
    if not os.path.isfile(idx): print("no global index yet — run a tagging pass first"); return
    runs=[];
    for ln in open(idx,encoding="utf-8",errors="replace"):
        try: runs.append(json.loads(ln))
        except Exception: pass
    tot=collections.Counter(); files=0
    for r in runs:
        files+=r.get("done",0)
        for k,c in (r.get("by") or {}).items(): tot[k]+=c
    print(f"docsort lifetime: {len(runs)} runs · {files} files tagged")
    for k,c in tot.most_common(25): print(f"  {k}: {c}")

def setup(a):
    global STREAMS,SUBJECTS,TYPES
    s,su,ty=load_tags(a.tags); STREAMS=set(s); SUBJECTS=set(su); TYPES=set(ty)
    return build_system(a.prompt,s,su,ty)

def add_args(ap):
    ap.add_argument("root",nargs="?",default=None,help="folder to process (or use --location)")
    ap.add_argument("--config",default=None,help="path to config.json (else the per-user config)")
    ap.add_argument("--host",default=None,help="override model host: a name from config 'hosts' or a raw URL")
    ap.add_argument("--list-models",action="store_true",help="list models loaded in LM Studio (at --host), then exit")
    ap.add_argument("--location",default=None,help="named location from config 'locations'")
    ap.add_argument("--api",default=None)
    ap.add_argument("--model",default=None)
    ap.add_argument("--vision",action="store_true",default=None)
    ap.add_argument("--vision-model",default=None)
    ap.add_argument("--backend",default=None,choices=["local"])
    ap.add_argument("--frontier",default=None,choices=["none","claude","cmd"])
    ap.add_argument("--frontier-cmd",dest="frontier_cmd",default=None,help="--frontier cmd: shell template; prompt piped via stdin")
    ap.add_argument("--tags",default=None,help="path to TAGS.md (default: your per-user copy)")
    ap.add_argument("--prompt",default=None,help="path to system_prompt.md (default: your per-user copy)")
    ap.add_argument("--edit-tags",action="store_true",help="open your TAGS.md in an editor, then exit")
    ap.add_argument("--apply",action="store_true",default=None)
    ap.add_argument("--copy",action="store_true",default=None,help="copy the folder to <name>COPY and tag the copy (originals untouched)")
    ap.add_argument("--misc",action=argparse.BooleanOptionalAction,default=True,
                    help="move 99UNS (unsure) files into a 'misc' subfolder (default ON; use --no-misc to disable)")
    ap.add_argument("--skip-unknown",dest="skip_unknown",action="store_true",
                    help="leave 99UNS (unknown) files completely untouched — no rename, no move")
    ap.add_argument("--move",default=None,help="destination root; or @archive to use config archive_root")
    ap.add_argument("--review",action="store_true",help="aggregate the run log into TAG-REVIEW.md (offline)")
    ap.add_argument("--report",action="store_true",help="(re)build DOCSORT-REPORT.md from the journal + update global index (offline)")
    ap.add_argument("--undo",action="store_true",help="reverse the renames/moves recorded in the journal")
    ap.add_argument("--stats",action="store_true",help="print lifetime stats from the global index, then exit")
    ap.add_argument("--retry-failed",dest="retry_failed",action="store_true",help="re-process only files marked 'failed' in the journal")
    ap.add_argument("--retag",action="store_true",help="re-classify already-prefixed files (after tuning/promoting a proposal)")
    ap.add_argument("--resume",action="store_true",help="skip files already done in the run journal (_docsort_state.jsonl)")
    ap.add_argument("--exclude",action="append",default=None,help="folder to skip (repeatable); adds to config 'exclude'")
    ap.add_argument("--include",action="append",default=None,help="only process this folder (repeatable); adds to config 'include'")
    ap.add_argument("--log",default=None)

def main(argv=None):
    ap=argparse.ArgumentParser(prog="docsort"); add_args(ap)
    pre,_=ap.parse_known_args(argv)
    if pre.edit_tags:                                # quick path: just open the tag list
        edit_file(tags_path()); return
    cfg=load_config(pre.config)
    ad,glb=arg_defaults(cfg); ap.set_defaults(**ad)
    a=ap.parse_args(argv)
    global MIN_TEXT,DEEP_PAGES,DEEP_CAP,DPI
    MIN_TEXT,DEEP_PAGES,DEEP_CAP,DPI=glb["MIN_TEXT"],glb["DEEP_PAGES"],glb["DEEP_CAP"],glb["DPI"]
    a.tags=a.tags or tags_path(); a.prompt=a.prompt or prompt_path()
    if a.list_models:                                # show what's loaded, then exit (no root needed)
        api=resolve_api(cfg,a.host) if a.host else a.api
        ms=available_models(api); print(f"host: {api}")
        if not ms: print("  (no response — is LM Studio running with a model loaded? check --host)")
        for m in ms:
            tag=" [VL]" if ("vl" in m.lower() or "vision" in m.lower()) else (" [embed]" if "embed" in m.lower() else "")
            print(f"  {m}{tag}")
        return
    if a.stats: stats(); return                              # global, no root needed
    root=resolve_location(cfg,a.location) if a.location else a.root
    if not root: ap.error("give a ROOT path or --location NAME (see config 'locations')")
    a.root=root
    if a.report: report(a.root); return                      # offline
    if a.undo: undo(a.root); return                          # offline
    if a.review: setup(a); review(a.root,a.log); return     # offline; no model server needed
    if a.move=="@archive":
        dest=cfg.get("archive_root")
        if not dest: ap.error("--move @archive but config 'archive_root' is empty")
    else: dest=a.move
    if dest: move_by_prefix(a.root,dest,bool(a.apply)); return
    if a.copy:                                   # safe mode: tag a copy, never the originals
        a.root=make_working_copy(a.root); print(f"[copy] working on {a.root}")
    if a.host: a.api=resolve_api(cfg,a.host)
    if a.frontier=="claude" and not shutil.which("claude"):
        print("[frontier] 'claude' CLI not on PATH — install Claude Code and run `claude` once to log in. "
              "Continuing without the frontier fallback."); a.frontier="none"
    if a.backend=="local":                       # adapt to whatever model is loaded
        m,up=resolve_model(a.api,a.model,prefer_vision=True)   # text tier too: grab any loaded VL model
        if not up: ap.error(f"model server unreachable at {a.api} — start LM Studio or fix --host (see TROUBLESHOOTING.md)")
        if m!=a.model: print(f"[model] '{a.model}' not loaded -> '{m}'"); a.model=m
        if a.vision:
            vm,up=resolve_model(a.api,a.vision_model,prefer_vision=True)
            if up and vm!=a.vision_model: print(f"[vision] '{a.vision_model}' not loaded -> '{vm}'"); a.vision_model=vm
    sysp=setup(a); rows=[]; n=0
    include=(cfg.get("include") or [])+(a.include or [])   # config + CLI
    exclude=(cfg.get("exclude") or [])+(a.exclude or [])
    items=[t for t in (iter_tagged(a.root) if a.retag else iter_targets(a.root))
           if passes_filter(os.path.join(t[0],t[1]),a.root,include,exclude)]
    if a.retry_failed:
        failed={j.get("rel") for j in _journal_rows(a.root) if j.get("status")=="failed"}
        items=[t for t in items if os.path.relpath(os.path.join(t[0],t[1]),a.root) in failed]
    if include or exclude: print(f"[filter] include={include or '-'} exclude={exclude or '-'} -> {len(items)} files")
    N=len(items)
    done_set=journal_done(a.root) if a.resume else {}
    jf=open(os.path.join(a.root,"_docsort_state.jsonl"),"a",encoding="utf-8")
    t0=time.time(); proc=0; dn=0; fl=0; toks=0; dead=0; skipped=0
    try:
        for dp,fn,base in items:
            full=os.path.join(dp,fn); rel=os.path.relpath(full,a.root)
            try: mt=int(os.path.getmtime(full))
            except Exception: mt=0
            if a.resume and done_set.get(rel)==mt:
                skipped+=1; continue                         # already done this exact file
            _STATS["ttoks"]=0; _STATS["ok"]=0; _STATS["empty"]=0
            status="done"; err=""
            try:
                st,su,ty,cf,src=classify(a,sysp,full,base,rel)
            except Exception as e:
                st,su,ty,cf,src=("CW","99UNS","misc","low","error"); status="failed"; err=str(e)[:200]
            toks+=_STATS["ttoks"]
            new=f"[{st}-{su}] {base}"; rows.append((full,fn,new,st,su,ty,cf,src)); n+=1; proc+=1
            if bool(getattr(a,"skip_unknown",False)) and su=="99UNS" and status=="done":
                status="skipped"; skipped+=1         # leave unknown files untouched (no rename/move)
            tomisc=bool(a.misc) and su=="99UNS" and status=="done"
            mark='  ->misc' if tomisc else ('  ->skip' if status=='skipped' else '')
            print(f"{st:5}{su:7}{ty:10}{cf:5}{src:14}{base[:46]}{mark}{'  FAIL' if status=='failed' else ''}")
            cur=full
            if a.apply and status=="done":
                if new!=fn:
                    nt=unique_path(os.path.join(dp,new))
                    try: os.rename(full,nt); cur=nt
                    except Exception as e: print("  rename-fail:",e); status="failed"; err="rename:"+str(e)[:160]
                if tomisc and status=="done":
                    try: cur=move_to_misc(a.root,cur)
                    except Exception as e: print("  misc-move-fail:",e)
            dn+=(status=="done"); fl+=(status=="failed")
            jf.write(json.dumps({"rel":rel,"name":base,"mtime":mt,"status":status,"stream":st,
                                 "subject":su,"type":ty,"conf":cf,"source":src,
                                 "dst":os.path.relpath(cur,a.root),"error":err,
                                 "ts":int(time.time())})+"\n"); jf.flush()
            el=time.time()-t0; tps=(toks/el) if el>0 else 0; avg=el/max(proc,1)
            eta=int(avg*max(N-proc-skipped,0))
            print(f"PROGRESS {proc}/{N} done={dn} failed={fl} tps={tps:.0f} toks={toks} eta={eta}s",flush=True)
            if status=="done" and _STATS["ok"]==0 and _STATS["empty"]>0: dead+=1
            else: dead=0
            if dead>=3 and not available_models(a.api):    # model server likely dropped mid-run
                jf.close(); _word_quit()
                print(f"\n[server] model unreachable at {a.api} — {dn} done. "
                      f"Fix LM Studio, then resume with: docsort \"{a.root}\" --resume")
                sys.exit(3)
    except KeyboardInterrupt:
        jf.close(); _word_quit()
        print(f"\n[paused] {dn} done, {fl} failed. Resume with: docsort \"{a.root}\" --resume")
        sys.exit(130)
    jf.close(); _word_quit()
    logp=a.log or os.path.join(a.root,"_docsort_log.csv")
    with open(logp,"w",encoding="utf-8",newline="") as g:
        w=csv.writer(g); w.writerow(["path","old","new","stream","subject","type","conf","source"]); w.writerows(rows)
    print(f"\n{'APPLIED' if a.apply else 'DRY-RUN'}: {n} files ({fl} failed). log: {logp}")
    report(a.root)                                           # DOCSORT-REPORT.md + global index
    if not a.apply: print("Review DOCSORT-REPORT.md (low-conf, failed, proposals), then --apply.")

if __name__=="__main__": main()
