#!/usr/bin/env python3
"""
doc_handler — tag academic documents with a local (or frontier) LLM, then move by prefix.

Tag vocabulary lives in TAGS.md (single source) and is injected into the model prompt.
Classification tiers: TEXT -> 5-page ESCALATE -> VISION (+page-3) -> FILENAME.
On a hard 99UNS, an optional FRONTIER model (Claude Code CLI / OpenAI API) gives a verdict.

  python doc_handler.py "D:\\AcademicsCOPY" --vision --vision-model qwen2-vl-7b
  python doc_handler.py "D:\\AcademicsCOPY" --apply --vision --vision-model qwen2-vl-7b
  python doc_handler.py "D:\\AcademicsCOPY" --frontier claude         # Claude Code sub for hard cases
  python doc_handler.py "D:\\AcademicsCOPY" --move "D:\\Archive\\Academics" --apply

Backends: --backend local|openai (main) ; --frontier none|claude|openai (99UNS fallback).
Deps: see requirements.txt.
"""
from __future__ import annotations
import os, sys, json, csv, re, base64, argparse, subprocess, shutil, urllib.request
from config import load_config, resolve_api, resolve_location, arg_defaults

def unique_path(p):
    """Avoid overwriting: foo.pdf -> foo__1.pdf if taken."""
    if not os.path.exists(p): return p
    base,ext=os.path.splitext(p); i=1
    while os.path.exists(f"{base}__{i}{ext}"): i+=1
    return f"{base}__{i}{ext}"

HERE=os.path.dirname(os.path.abspath(__file__))
EXT_TEXT={".pdf",".txt",".md",".docx",".pptx",".ppt",".doc"}
MIN_TEXT=80; DEEP_PAGES=5; DEEP_CAP=4000; DPI=120
STREAMS=set(); SUBJECTS=set(); TYPES=set()   # filled from TAGS.md

# ---------- tag loading (single source) ----------
def load_tags(path):
    streams={}; subjects={}; types=[]
    try: txt=open(path,encoding="utf-8").read()
    except Exception: txt=""
    def block(header):
        m=re.search(r"##\s*"+header+r".*?```tags\n(.*?)```",txt,re.S|re.I)
        return m.group(1).strip().splitlines() if m else []
    for ln in block("STREAMS"):
        if ln.strip(): streams[ln.split()[0]]=ln.strip()
    for ln in block("SUBJECTS"):
        if ln.strip(): subjects[ln.split()[0]]=ln.strip()
    m=re.search(r"##\s*TYPES.*?```tags\n(.*?)```",txt,re.S|re.I)
    if m: types=[x.split()[0] for x in m.group(1).splitlines() if x.strip()]
    if not streams: streams={"CW":"CW","GATE":"GATE","PROJ":"PROJ","RES":"RES","REC":"REC","REF":"REF"}
    if not subjects: subjects={"99UNS":"unsure","NA":"na"}
    if not types: types=["misc"]
    return streams,subjects,types

def build_system(template_path, streams, subjects, types):
    try:
        s=open(template_path,encoding="utf-8").read()
        m=re.search(r"<SYSTEM>(.*?)</SYSTEM>",s,re.S); tmpl=m.group(1) if m else s
    except Exception:
        tmpl="Reply: STREAM SUBJECT TYPE CONF.\n{{STREAMS}}\n{{SUBJECTS}}\n{{TYPES}}"
    tmpl=tmpl.replace("{{STREAMS}}","\n".join("  "+v for v in streams.values()))
    tmpl=tmpl.replace("{{SUBJECTS}}","\n".join("  "+v for v in subjects.values()))
    tmpl=tmpl.replace("{{TYPES}}"," ".join(types))
    return tmpl.strip()

# ---------- text / image extraction ----------
def pdf_text(path,pages=2,cap=2000):
    try:
        import fitz; d=fitz.open(path); t=""
        for i in range(min(pages,d.page_count)):
            t+=d[i].get_text()
            if len(t)>cap: break
        d.close(); return t[:cap]
    except Exception: return ""
_WORD=None
def _doc_legacy(path):
    """Read old binary .doc via a reused Word COM instance (Windows + Word + pywin32)."""
    global _WORD
    try:
        if _WORD is None:
            import win32com.client as wc
            _WORD=wc.Dispatch("Word.Application"); _WORD.Visible=False
        d=_WORD.Documents.Open(os.path.abspath(path), ReadOnly=True, AddToRecentFiles=False)
        t=d.Content.Text; d.Close(False); return t
    except Exception: return ""   # no Word/pywin32 -> falls back to filename tier
def _word_quit():
    global _WORD
    try:
        if _WORD is not None: _WORD.Quit(); _WORD=None
    except Exception: pass

def doc_text(path,cap=2000):
    e=os.path.splitext(path)[1].lower()
    try:
        if e==".pdf": return pdf_text(path,2,cap)
        if e in (".txt",".md"): return open(path,encoding="utf-8",errors="replace").read()[:cap]
        if e==".doc": return _doc_legacy(path)[:cap]
        if e==".docx":
            import docx; return "\n".join(p.text for p in docx.Document(path).paragraphs)[:cap]
        if e==".pptx":
            from pptx import Presentation
            return " ".join(s.text for sl in Presentation(path).slides for s in sl.shapes
                             if getattr(s,'has_text_frame',False))[:cap]
    except Exception: return ""
    return ""
def page_png(path,page=0,dpi=None):
    try:
        import fitz; d=fitz.open(path)
        if d.page_count==0: return None
        if d.page_count<=page: page=0
        b=d[page].get_pixmap(dpi=dpi or DPI).tobytes("png"); d.close(); return b
    except Exception: return None

# ---------- backends ----------
def _http(url,payload,headers):
    req=urllib.request.Request(url,data=json.dumps(payload).encode(),headers=headers)
    return json.load(urllib.request.urlopen(req,timeout=180))

def available_models(api):
    try:
        r=json.load(urllib.request.urlopen(api.replace("/chat/completions","/models"),timeout=8))
        return [m.get("id") for m in r.get("data",[])]
    except Exception: return []

def resolve_model(api,want,prefer_vision=False):
    """If 'want' isn't loaded on the server, fall back to a loaded model (vision-preferring).
    Returns (model_id, server_reachable)."""
    av=available_models(api)
    if not av: return want, False
    if want in av: return want, True
    pick=None
    if prefer_vision:
        pick=next((m for m in av if ("vl" in m.lower() or "vision" in m.lower())),None)
    pick=pick or next((m for m in av if "embed" not in m.lower()), av[0])
    return pick, True

def llm(a,backend,system,user_text,image_png=None):
    """Return raw model text. backend: local | openai | claude."""
    if backend in ("local","openai"):
        if image_png:
            content=[{"type":"text","text":user_text},
                     {"type":"image_url","image_url":{"url":"data:image/png;base64,"+base64.b64encode(image_png).decode()}}]
        else:
            content=user_text
        msgs=[{"role":"system","content":system},{"role":"user","content":content}]
        try:
            if backend=="local":
                model=a.vision_model if image_png else a.model
                r=_http(a.api,{"model":model,"temperature":0,"max_tokens":24,"messages":msgs},
                        {"Content-Type":"application/json"})
            else:
                key=os.environ.get("OPENAI_API_KEY","")
                r=_http("https://api.openai.com/v1/chat/completions",
                        {"model":a.openai_model,"temperature":0,"max_tokens":24,"messages":msgs},
                        {"Content-Type":"application/json","Authorization":"Bearer "+key})
            return r["choices"][0]["message"]["content"]
        except Exception as e:
            return ""   # -> decide() -> 99UNS; one bad call never kills the run
    if backend in ("claude","cmd"):   # subscription CLIs — NO API key. Text only.
        import tempfile
        oneshot=(system+"\n\n"+user_text+
                 "\n\nRespond with ONLY one line: STREAM SUBJECT TYPE CONF. "
                 "No explanation, no markdown, no questions.")
        neutral=tempfile.gettempdir()   # run outside the repo so the CLI doesn't read project files
        try:
            if backend=="claude":       # Claude Code CLI -> Claude subscription
                exe=shutil.which("claude") or "claude"
                fm=getattr(a,"frontier_model","") or "haiku"   # standard-context, sub-covered
                out=subprocess.run([exe,"-p","--model",fm,oneshot],capture_output=True,text=True,
                                   timeout=180,cwd=neutral)   # resolved .exe -> no shell (shell+list mangles on Win)
                return out.stdout
            tmpl=getattr(a,"frontier_cmd","") or ""   # cmd: templated CLI from config, prompt via stdin
            if not tmpl: return ""
            import shlex
            out=subprocess.run(shlex.split(tmpl),input=oneshot,capture_output=True,text=True,
                               timeout=180,cwd=neutral)
            return out.stdout
        except Exception: return ""
    return ""

def parse(out):
    o=(out or "").upper()
    st=next((s for s in sorted(STREAMS,key=len,reverse=True) if re.search(r'\b'+re.escape(s)+r'\b',o)),"CW")
    su=next((c for c in sorted(SUBJECTS,key=len,reverse=True) if re.search(r'\b'+re.escape(c)+r'\b',o)),"99UNS")
    ty=next((t for t in TYPES if re.search(r'\b'+re.escape(t.upper())+r'\b',o)),"misc")
    cf="high" if "HIGH" in o else "low"
    m=re.search(r'PROPOSE[:\s]+([A-Z][A-Z0-9_]{1,11})',o)
    prop=m.group(1) if m else ""
    return st,su,ty,cf,prop

def decide(out):
    """parse + fold a model proposal into a ~LABEL subject (review symbol)."""
    st,su,ty,cf,prop=parse(out)
    if prop and su=="99UNS": su="~"+prop
    return st,su,ty,cf

def classify(a,sysp,full,fn,rel):
    ispdf=full.lower().endswith(".pdf")
    snip=doc_text(full)
    u=lambda txt:f"Filename: {fn}\nFolder: {rel}\nText:\n{txt[:DEEP_CAP]}\n\nAnswer (STREAM SUBJECT TYPE CONF):"
    if len(snip.strip())>=MIN_TEXT:
        st,su,ty,cf=decide(llm(a,a.backend,sysp,u(snip))); src="text"
        if su=="99UNS" and ispdf:
            deep=pdf_text(full,DEEP_PAGES,DEEP_CAP)
            if len(deep)>len(snip):
                r=decide(llm(a,a.backend,sysp,u(deep)))
                if r[1]!="99UNS": st,su,ty,cf,src=(*r,"text5"); snip=deep
        if su=="99UNS" and a.frontier!="none":
            r=decide(llm(a,a.frontier,sysp,u(snip)))
            if r[1]!="99UNS": return (*r,"frontier:"+a.frontier)
        return st,su,ty,cf,src
    if a.vision and ispdf and (png:=page_png(full,0)):
        un="(handwritten/scanned page) Answer (STREAM SUBJECT TYPE CONF):"
        st,su,ty,cf=decide(llm(a,a.backend,sysp,f"Filename: {fn}\nFolder: {rel}\n"+un,png)); src="vision"
        if su=="99UNS" and (p3:=page_png(full,2)):
            r=decide(llm(a,a.backend,sysp,f"Filename: {fn}\nFolder: {rel}\n(page 3) "+un,p3))
            if r[1]!="99UNS": return (*r,"vision3")
        return st,su,ty,cf,src
    r=decide(llm(a,a.backend,sysp,u("")))
    if r[1]=="99UNS" and a.frontier!="none":
        rf=decide(llm(a,a.frontier,sysp,u("")))
        if rf[1]!="99UNS": return (*rf,"frontier:"+a.frontier)
    return (*r,"filename")

def move_by_prefix(root,dest,apply):
    pat=re.compile(r'^\[([A-Z]+)-([0-9A-Z]+)\]\s+(.*)$'); n=0; rows=[]
    for dp,_,fns in os.walk(root):
        for fn in fns:
            m=pat.match(fn)
            if not m: continue
            src=os.path.join(dp,fn); tgt=os.path.join(dest,m.group(1),m.group(2))
            newp=os.path.join(tgt,m.group(3))
            if apply:
                os.makedirs(tgt,exist_ok=True)
                newp=unique_path(newp)                       # never overwrite
                try: shutil.move(src,newp)                   # cross-device safe
                except Exception as e: print("move-fail:",fn,e)
            rows.append((src,newp)); n+=1
    log=os.path.join(dest if apply else root,"_move_log.csv")
    try:
        os.makedirs(os.path.dirname(log),exist_ok=True)
        with open(log,"w",encoding="utf-8",newline="") as g: csv.writer(g).writerows([("from","to")]+rows)
    except Exception: pass
    print(f"{'MOVED' if apply else 'DRY-RUN move'}: {n} files -> {dest}")
    return n

SKIP_NAMES={"tag-review.md","readme.md","changelog.md","guide.md","troubleshooting.md","tags.md"}
def iter_targets(root):
    for dp,_,fns in os.walk(root):
        for fn in fns:
            if fn.lower() in SKIP_NAMES or fn.startswith(("_doc_handler","_move")): continue
            if os.path.splitext(fn)[1].lower() in EXT_TEXT and not fn.startswith("["):
                yield dp,fn,fn                       # (dir, current name, base name)

_TAGPFX=re.compile(r'^\[[^\]]+\]\s+(.*)$')
def iter_tagged(root):
    """Already-prefixed docs, for --retag. Yields (dir, current name, stripped base)."""
    for dp,_,fns in os.walk(root):
        for fn in fns:
            m=_TAGPFX.match(fn)
            if m and os.path.splitext(m.group(1))[1].lower() in EXT_TEXT:
                yield dp,fn,m.group(1)

def review(root, log=None):
    """Aggregate a run log into TAG-REVIEW.md: distribution, proposed (~) tags, low-conf. No model needed."""
    import collections
    logp=log or os.path.join(root,"_doc_handler_log.csv")
    if not os.path.isfile(logp): print("no log at",logp); return
    rows=list(csv.DictReader(open(logp,encoding="utf-8")))
    combo=collections.Counter(f"{r['stream']}-{r['subject']}" for r in rows)
    props=collections.Counter(r['subject'] for r in rows if r['subject'].startswith('~'))
    low=[r for r in rows if r['conf']=='low']
    out=["# TAG REVIEW","",f"{len(rows)} files · {len(props)} proposed tag(s) · {len(low)} low-confidence.","",
         "## Distribution (stream-subject)"]
    for k,n in combo.most_common(): out.append(f"- `{k}` x {n}")
    out.append("\n## (review) Proposed tags -> promote to TAGS.md")
    if props:
        for lab,n in props.most_common():
            out.append(f"- `{lab}` x {n}  -> if recurring, add a SUBJECT code for `{lab[1:]}` in TAGS.md, then re-run")
    else: out.append("- none")
    out.append("\n## Low-confidence (eyeball)")
    for r in low[:50]: out.append(f"- {r['stream']}-{r['subject']} · {r['source']} · {r['old']}")
    p=os.path.join(root,"TAG-REVIEW.md"); open(p,"w",encoding="utf-8").write("\n".join(out))
    print(f"review -> {p}")
    print(f"proposals: {dict(props) or 'none'} | low-conf: {len(low)} | groups: {len(combo)}")

def setup(a):
    global STREAMS,SUBJECTS,TYPES
    s,su,ty=load_tags(a.tags); STREAMS=set(s); SUBJECTS=set(su); TYPES=set(ty)
    return build_system(a.prompt,s,su,ty)

def add_args(ap):
    ap.add_argument("root",nargs="?",default=None,help="folder to process (or use --location)")
    ap.add_argument("--config",default=None,help="path to config.json (else config.json beside script)")
    ap.add_argument("--host",default=None,help="override model host: a name from config 'hosts' or a raw URL")
    ap.add_argument("--location",default=None,help="named location from config 'locations'")
    ap.add_argument("--api",default=None)
    ap.add_argument("--model",default=None)
    ap.add_argument("--vision",action="store_true",default=None)
    ap.add_argument("--vision-model",default=None)
    ap.add_argument("--backend",default=None,choices=["local","openai"])
    ap.add_argument("--frontier",default=None,choices=["none","claude","openai","cmd"])
    ap.add_argument("--frontier-model",default="haiku",help="Claude Code model for --frontier claude (standard-context, sub-covered)")
    ap.add_argument("--frontier-cmd",dest="frontier_cmd",default=None,help="--frontier cmd: shell template; prompt piped via stdin")
    ap.add_argument("--openai-model",default=None)
    ap.add_argument("--tags",default=os.path.join(HERE,"TAGS.md"))
    ap.add_argument("--prompt",default=os.path.join(HERE,"system_prompt.md"))
    ap.add_argument("--apply",action="store_true",default=None)
    ap.add_argument("--move",default=None,help="destination root; or @archive to use config archive_root")
    ap.add_argument("--review",action="store_true",help="aggregate the run log into TAG-REVIEW.md (offline)")
    ap.add_argument("--retag",action="store_true",help="re-classify already-prefixed files (after tuning/promoting a proposal)")
    ap.add_argument("--log",default=None)

def main():
    ap=argparse.ArgumentParser(); add_args(ap)
    pre,_=ap.parse_known_args()
    cfg=load_config(pre.config)
    ad,glb=arg_defaults(cfg); ap.set_defaults(**ad)
    a=ap.parse_args()
    global MIN_TEXT,DEEP_PAGES,DEEP_CAP,DPI
    MIN_TEXT,DEEP_PAGES,DEEP_CAP,DPI=glb["MIN_TEXT"],glb["DEEP_PAGES"],glb["DEEP_CAP"],glb["DPI"]
    root=resolve_location(cfg,a.location) if a.location else a.root
    if not root: ap.error("give a ROOT path or --location NAME (see config 'locations')")
    a.root=root
    if a.review: setup(a); review(a.root,a.log); return     # offline; no model server needed
    if a.move=="@archive":
        dest=cfg.get("archive_root")
        if not dest: ap.error("--move @archive but config 'archive_root' is empty")
    else: dest=a.move
    if dest: move_by_prefix(a.root,dest,bool(a.apply)); return
    if a.host: a.api=resolve_api(cfg,a.host)
    if a.backend=="local":                       # adapt to whatever model is loaded
        m,up=resolve_model(a.api,a.model)
        if not up: ap.error(f"model server unreachable at {a.api} — start LM Studio or fix --host (see TROUBLESHOOTING.md)")
        if m!=a.model: print(f"[model] '{a.model}' not loaded -> '{m}'"); a.model=m
        if a.vision:
            vm,up=resolve_model(a.api,a.vision_model,prefer_vision=True)
            if up and vm!=a.vision_model: print(f"[vision] '{a.vision_model}' not loaded -> '{vm}'"); a.vision_model=vm
    sysp=setup(a); rows=[]; n=0
    items=iter_tagged(a.root) if a.retag else iter_targets(a.root)
    for dp,fn,base in items:
        full=os.path.join(dp,fn); rel=os.path.relpath(dp,a.root)
        st,su,ty,cf,src=classify(a,sysp,full,base,rel)
        new=f"[{st}-{su}] {base}"; rows.append((full,fn,new,st,su,ty,cf,src)); n+=1
        print(f"{st:5}{su:7}{ty:10}{cf:5}{src:14}{base[:46]}")
        if a.apply and new!=fn:
            try: os.rename(full,unique_path(os.path.join(dp,new)))
            except Exception as e: print("  rename-fail:",e)
    _word_quit()
    logp=a.log or os.path.join(a.root,"_doc_handler_log.csv")
    with open(logp,"w",encoding="utf-8",newline="") as g:
        w=csv.writer(g); w.writerow(["path","old","new","stream","subject","type","conf","source"]); w.writerows(rows)
    print(f"\n{'APPLIED' if a.apply else 'DRY-RUN'}: {n} files. log: {logp}")
    if not a.apply: print("Review log (conf=low, source=filename/vision/frontier), then --apply.")

if __name__=="__main__": main()
